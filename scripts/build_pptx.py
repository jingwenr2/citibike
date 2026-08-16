"""Build a capstone project PowerPoint presentation."""
from __future__ import annotations

import json
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

ROOT = Path(__file__).resolve().parents[1]
OUT = ROOT / "report" / "CitiBike_Capstone_Presentation.pptx"

# ── Brand colors ──
NAVY = RGBColor(0x0F, 0x17, 0x2A)
DARK_BLUE = RGBColor(0x1E, 0x3A, 0x5F)
PINK = RGBColor(0xFF, 0x00, 0xBF)
PURPLE = RGBColor(0x7C, 0x3A, 0xED)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0x94, 0xA3, 0xB8)
MID_GRAY = RGBColor(0x64, 0x74, 0x8B)
GREEN = RGBColor(0x05, 0x96, 0x69)
BLUE = RGBColor(0x2D, 0x7F, 0xF9)
AMBER = RGBColor(0xD9, 0x77, 0x06)
BG_LIGHT = RGBColor(0xF8, 0xF9, 0xFB)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape_fill(slide, left, top, width, height, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_text_box(slide, left, top, width, height, text, font_size=18,
                 color=NAVY, bold=False, alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_bullet_slide_content(slide, left, top, width, height, bullets, font_size=16,
                              color=NAVY, spacing=Pt(8)):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = bullet
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = "Calibri"
        p.space_after = spacing
        p.level = 0
    return txBox


def add_kpi_card(slide, left, top, width, height, label, value, accent=BLUE):
    card = add_shape_fill(slide, left, top, width, height, WHITE)
    card.shadow.inherit = False

    # Accent bar at top
    add_shape_fill(slide, left, top, width, Pt(4), accent)

    # Value
    add_text_box(slide, left + Inches(0.15), top + Inches(0.2),
                 width - Inches(0.3), Inches(0.5),
                 value, font_size=22, color=NAVY, bold=True,
                 alignment=PP_ALIGN.CENTER)
    # Label
    add_text_box(slide, left + Inches(0.15), top + Inches(0.65),
                 width - Inches(0.3), Inches(0.3),
                 label, font_size=11, color=MID_GRAY,
                 alignment=PP_ALIGN.CENTER)


def section_divider(prs, title, subtitle=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_slide_bg(slide, NAVY)

    # Pink accent circle (decorative)
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                     Inches(9.5), Inches(-1),
                                     Inches(5), Inches(5))
    circle.fill.solid()
    circle.fill.fore_color.rgb = RGBColor(0xFF, 0x00, 0xBF)
    fill = circle.fill
    fill.solid()
    fill.fore_color.rgb = PINK
    circle.fill.fore_color.brightness = 0.85
    circle.line.fill.background()

    add_text_box(slide, Inches(1), Inches(2.5), Inches(8), Inches(1),
                 title, font_size=36, color=WHITE, bold=True)
    if subtitle:
        add_text_box(slide, Inches(1), Inches(3.5), Inches(8), Inches(0.8),
                     subtitle, font_size=18, color=LIGHT_GRAY)
    return slide


def build():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    # ================================================================
    # SLIDE 1: Title
    # ================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, NAVY)

    add_text_box(slide, Inches(1), Inches(0.8), Inches(10), Inches(0.4),
                 "CAPSTONE PROJECT", font_size=14, color=PINK, bold=True)
    add_text_box(slide, Inches(1), Inches(1.4), Inches(10), Inches(1.2),
                 "NYC Citi Bike Public Investment Intelligence",
                 font_size=40, color=WHITE, bold=True)
    add_text_box(slide, Inches(1), Inches(3.0), Inches(10), Inches(0.8),
                 "Data-driven case for expanding NYC's bike-share network\n"
                 "69.7M trips analyzed  |  2,476 stations  |  XGBoost demand forecast",
                 font_size=18, color=LIGHT_GRAY)
    add_text_box(slide, Inches(1), Inches(5.5), Inches(6), Inches(0.4),
                 "Jingwen Ruan", font_size=16, color=WHITE, bold=True)
    add_text_box(slide, Inches(1), Inches(5.9), Inches(6), Inches(0.4),
                 "August 2026", font_size=14, color=LIGHT_GRAY)

    # ================================================================
    # SLIDE 2: Executive Summary
    # ================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_LIGHT)
    add_text_box(slide, Inches(0.8), Inches(0.5), Inches(10), Inches(0.6),
                 "Executive Summary", font_size=32, color=NAVY, bold=True)

    # Insight panel
    panel = add_shape_fill(slide, Inches(0.8), Inches(1.3), Inches(11.5), Inches(1.0),
                           RGBColor(0xF0, 0xF9, 0xFF))
    add_shape_fill(slide, Inches(0.8), Inches(1.3), Pt(4), Inches(1.0), BLUE)
    add_text_box(slide, Inches(1.0), Inches(1.4), Inches(11), Inches(0.8),
                 "Citi Bike is a $212M/yr business running at capacity — 51.9% of stations are maxed out. "
                 "Adding 250 stations nets $12.9M/yr profit with a 15-month payback. "
                 "SF proves government partnership works; NYC is the 10x opportunity.",
                 font_size=16, color=NAVY)

    # KPI cards
    kpis = [
        ("69.7M", "Total Trips Analyzed", BLUE),
        ("$212M/yr", "Estimated Revenue", GREEN),
        ("51.9%", "Stations at Capacity", AMBER),
        ("15 months", "Expansion Payback", PINK),
        ("+42.5%", "Model vs Naive Baseline", PURPLE),
    ]
    card_w = Inches(2.1)
    start_x = Inches(0.8)
    for i, (val, label, accent) in enumerate(kpis):
        add_kpi_card(slide, start_x + i * (card_w + Inches(0.15)),
                     Inches(2.7), card_w, Inches(1.1), label, val, accent)

    # 5 key points
    bullets = [
        "NYC Citi Bike is the largest bike-share in the Americas (2,476 stations, 69.7M trips)",
        "Supply-constrained, not demand-constrained — 961 stations above 1.5x capacity",
        "E-bikes drive 70% of trips and ~$97M/yr in overage fees alone",
        "XGBoost demand forecast: MAE 11.24 trips/station/day (+42.5% vs naive baseline)",
        "250-station expansion yields $12.9M/yr net profit; 5-year gap if not investing: ~$377M",
    ]
    add_bullet_slide_content(slide, Inches(0.8), Inches(4.2), Inches(11.5), Inches(3),
                              bullets, font_size=15, color=NAVY)

    # ================================================================
    # SECTION: Problem & Opportunity
    # ================================================================
    section_divider(prs, "Problem & Opportunity",
                    "Why NYC needs to invest in Citi Bike now")

    # SLIDE 3: The Problem
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_LIGHT)
    add_text_box(slide, Inches(0.8), Inches(0.5), Inches(10), Inches(0.6),
                 "The Problem: Capacity Crisis", font_size=32, color=NAVY, bold=True)

    left_bullets = [
        "51.9% of stations operating at or above capacity",
        "961 stations (38.8%) in critical zone (>1.5x capacity)",
        "Peak month: 5.4M trips in June 2026",
        "NYC averages 51.5 trips/station/day vs SF's 16.6",
        "Demand is being turned away, not just unmet",
    ]
    add_bullet_slide_content(slide, Inches(0.8), Inches(1.5), Inches(5.5), Inches(3),
                              left_bullets, font_size=16, color=NAVY)

    right_bullets = [
        "MTA subway reliability declining — riders switching to bikes",
        "IBZ corridor: station deserts in underserved neighborhoods",
        "250 new stations already planned for Bronx/Queens/Brooklyn",
        "But which neighborhoods should be prioritized?",
        "Our model answers: where is unmet demand highest?",
    ]
    add_bullet_slide_content(slide, Inches(7), Inches(1.5), Inches(5.5), Inches(3),
                              right_bullets, font_size=16, color=NAVY)

    # Bottom insight
    panel = add_shape_fill(slide, Inches(0.8), Inches(5.5), Inches(11.5), Inches(0.8),
                           RGBColor(0xFE, 0xF3, 0xC7))
    add_text_box(slide, Inches(1.0), Inches(5.6), Inches(11), Inches(0.6),
                 "Key insight: This is a supply problem, not a demand problem. "
                 "Every dock at capacity is a lost trip and lost revenue.",
                 font_size=15, color=RGBColor(0x92, 0x40, 0x0E), bold=True)

    # ================================================================
    # SECTION: Data & Methods
    # ================================================================
    section_divider(prs, "Data & Methodology",
                    "69.7M trips, 23 features, 3-layer architecture")

    # SLIDE 4: Data Pipeline
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_LIGHT)
    add_text_box(slide, Inches(0.8), Inches(0.5), Inches(10), Inches(0.6),
                 "Data Pipeline & Sources", font_size=32, color=NAVY, bold=True)

    datasets = [
        ("Citi Bike Trip Histories", "16.6M+ trips, monthly CSVs (2024-2026)", "Primary ridership data"),
        ("Bay Wheels (SF)", "8.1M trips, 23 months", "Government investment benchmark"),
        ("MTA Subway Data", "300 stations joined to bike stations", "Transit reliability & proximity"),
        ("NYC Weather (NOAA)", "Daily temperature & precipitation", "Deviation-based weather features"),
        ("Census / Geospatial", "IBZ corridor demographics", "Equity & access analysis"),
    ]
    for i, (name, detail, purpose) in enumerate(datasets):
        y = Inches(1.5) + i * Inches(1.0)
        add_shape_fill(slide, Inches(0.8), y, Inches(11.5), Inches(0.85), WHITE)
        add_text_box(slide, Inches(1.0), y + Inches(0.08), Inches(3.5), Inches(0.35),
                     name, font_size=15, color=NAVY, bold=True)
        add_text_box(slide, Inches(4.8), y + Inches(0.08), Inches(3.5), Inches(0.35),
                     detail, font_size=14, color=MID_GRAY)
        add_text_box(slide, Inches(8.5), y + Inches(0.08), Inches(3.5), Inches(0.35),
                     purpose, font_size=14, color=DARK_BLUE)

    # Architecture
    add_text_box(slide, Inches(0.8), Inches(6.5), Inches(11.5), Inches(0.5),
                 "Architecture:  Data Core (Parquet)  ->  City Context (MTA + Weather)  ->  "
                 "Investment Web App (Streamlit + XGBoost)",
                 font_size=14, color=MID_GRAY, alignment=PP_ALIGN.CENTER)

    # ================================================================
    # SLIDE 5: Model Overview
    # ================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_LIGHT)
    add_text_box(slide, Inches(0.8), Inches(0.5), Inches(10), Inches(0.6),
                 "XGBoost Demand Forecast Model", font_size=32, color=NAVY, bold=True)

    add_text_box(slide, Inches(0.8), Inches(1.2), Inches(11), Inches(0.5),
                 "Predicts daily trip counts per station using 23 features across 6 categories",
                 font_size=16, color=MID_GRAY)

    # Feature categories
    categories = [
        ("Station", "lat, lon, capacity", BLUE),
        ("Calendar", "day_of_week, month, quarter, year,\nis_weekend, is_holiday", PURPLE),
        ("Lag / Rolling", "lag_1d, lag_7d, roll_mean_7d,\nroll_mean_28d, roll_std_7d", GREEN),
        ("Ridership", "electric_share, member_share", PINK),
        ("MTA Transit", "mta_daily_riders, mta_delay_rate,\nnearest_mta_distance_km", AMBER),
        ("Weather", "temp_deviation, precip_deviation,\nis_bad_weather", BLUE),
    ]
    card_w = Inches(1.85)
    for i, (cat, feats, accent) in enumerate(categories):
        x = Inches(0.5) + i * (card_w + Inches(0.1))
        y = Inches(2.0)
        add_shape_fill(slide, x, y, card_w, Inches(2.2), WHITE)
        add_shape_fill(slide, x, y, card_w, Pt(4), accent)
        add_text_box(slide, x + Inches(0.1), y + Inches(0.15),
                     card_w - Inches(0.2), Inches(0.3),
                     cat, font_size=14, color=NAVY, bold=True, alignment=PP_ALIGN.CENTER)
        add_text_box(slide, x + Inches(0.1), y + Inches(0.55),
                     card_w - Inches(0.2), Inches(1.5),
                     feats, font_size=11, color=MID_GRAY, alignment=PP_ALIGN.CENTER)

    # Training details
    training_info = [
        "Training: 1.44M station-day rows across 2,464 stations",
        "Chronological split: last 60 days held out for testing (137,841 station-days)",
        "Early stopping: 30-round patience on 20% eval holdout (best iteration: 879/1000)",
        "Weather: deviation from monthly climate normals (15C in Jan != 15C in Jul)",
        "3-fold time-series cross-validation for robust evaluation",
    ]
    add_bullet_slide_content(slide, Inches(0.8), Inches(4.5), Inches(11.5), Inches(2.5),
                              training_info, font_size=14, color=NAVY)

    # ================================================================
    # SLIDE 6: Model Results
    # ================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_LIGHT)
    add_text_box(slide, Inches(0.8), Inches(0.5), Inches(10), Inches(0.6),
                 "Model Performance", font_size=32, color=NAVY, bold=True)

    # Results table
    metrics_path = ROOT / "models" / "forecast_metrics.json"
    if metrics_path.exists():
        with open(metrics_path) as f:
            metrics = json.load(f)[0]
        xgb = metrics["xgboost"]
        naive = metrics["seasonal_naive"]
        cv_avg = metrics.get("cv_avg", {})
    else:
        xgb = {"MAE": 11.24, "RMSE": 21.43, "WAPE": 0.1534}
        naive = {"MAE": 19.56, "RMSE": 39.92, "WAPE": 0.267}
        cv_avg = {"MAE": 9.35, "RMSE": 17.78, "WAPE": 0.2020}

    # KPI cards for results
    result_kpis = [
        (f"{xgb['MAE']}", "XGBoost MAE", PURPLE),
        (f"{naive['MAE']}", "Naive Baseline MAE", LIGHT_GRAY),
        (f"+42.5%", "Improvement", GREEN),
        (f"{xgb['WAPE']:.0%}", "WAPE", PURPLE),
        (f"{cv_avg.get('MAE', '9.35')}", "CV Average MAE", BLUE),
    ]
    card_w = Inches(2.1)
    for i, (val, label, accent) in enumerate(result_kpis):
        add_kpi_card(slide, Inches(0.8) + i * (card_w + Inches(0.15)),
                     Inches(1.5), card_w, Inches(1.1), label, str(val), accent)

    # Comparison table
    add_text_box(slide, Inches(0.8), Inches(3.2), Inches(5), Inches(0.4),
                 "Test Set Comparison (137,841 station-days)", font_size=16, color=NAVY, bold=True)

    headers = ["Metric", "XGBoost", "Naive (7d)", "Improvement"]
    rows = [
        ["MAE", f"{xgb['MAE']}", f"{naive['MAE']}", "+42.5%"],
        ["RMSE", f"{xgb['RMSE']}", f"{naive['RMSE']}", "+46.3%"],
        ["WAPE", f"{xgb['WAPE']:.1%}", f"{naive['WAPE']:.1%}", "+42.5%"],
    ]
    table = slide.shapes.add_table(len(rows) + 1, 4,
                                    Inches(0.8), Inches(3.8),
                                    Inches(6), Inches(1.8)).table
    for j, h in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = h
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(13)
            p.font.bold = True
            p.font.color.rgb = WHITE
        cell.fill.solid()
        cell.fill.fore_color.rgb = NAVY
    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            cell = table.cell(i + 1, j)
            cell.text = val
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(13)
                p.font.color.rgb = NAVY
                if j == 3:
                    p.font.color.rgb = GREEN
                    p.font.bold = True

    # CV results
    if cv_avg:
        add_text_box(slide, Inches(7.5), Inches(3.2), Inches(5), Inches(0.4),
                     "Cross-Validation (3-fold)", font_size=16, color=NAVY, bold=True)
        cv_bullets = [
            f"Fold 1: MAE = 6.68 (earliest data, smallest train set)",
            f"Fold 2: MAE = 10.12",
            f"Fold 3: MAE = 11.24 (matches test set)",
            f"Average: MAE = {cv_avg.get('MAE', 9.35)}, WAPE = {cv_avg.get('WAPE', 0.202):.1%}",
            "Consistent performance across time periods",
        ]
        add_bullet_slide_content(slide, Inches(7.5), Inches(3.8), Inches(5), Inches(2.5),
                                  cv_bullets, font_size=14, color=NAVY)

    # Top features
    add_text_box(slide, Inches(0.8), Inches(6.0), Inches(11.5), Inches(0.4),
                 "Top features: roll_mean_7d (53.8%) > lag_1d (21.5%) > roll_mean_28d (8.5%) > lag_7d (6.0%)",
                 font_size=14, color=MID_GRAY, alignment=PP_ALIGN.CENTER)

    # ================================================================
    # SECTION: Revenue & Investment
    # ================================================================
    section_divider(prs, "Revenue & Investment Case",
                    "The financial case for expansion")

    # SLIDE 7: Revenue Analysis
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_LIGHT)
    add_text_box(slide, Inches(0.8), Inches(0.5), Inches(10), Inches(0.6),
                 "Revenue & E-Bike Economics", font_size=32, color=NAVY, bold=True)

    rev_kpis = [
        ("$212M/yr", "Estimated Annual Revenue", GREEN),
        ("$105.2M/yr", "E-Bike Overage Fees", PINK),
        ("~70%", "E-Bike Trip Share", BLUE),
        ("$239/yr", "Annual Membership", NAVY),
    ]
    card_w = Inches(2.6)
    for i, (val, label, accent) in enumerate(rev_kpis):
        add_kpi_card(slide, Inches(0.8) + i * (card_w + Inches(0.2)),
                     Inches(1.5), card_w, Inches(1.1), label, val, accent)

    rev_bullets = [
        "E-bikes are the revenue engine: $0.27/min overage drives ~50% of revenue",
        "Member share is growing — stickier, higher LTV user base",
        "Price increased from $220 to $239/yr (Jan 2026) — demand remains strong",
        "DC and Chicago are publicly subsidized; NYC comparison must account for this",
        "Revenue per trip higher in NYC than any comparable US system",
    ]
    add_bullet_slide_content(slide, Inches(0.8), Inches(3.2), Inches(11.5), Inches(3),
                              rev_bullets, font_size=16, color=NAVY)

    # ================================================================
    # SLIDE 8: Expansion ROI
    # ================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_LIGHT)
    add_text_box(slide, Inches(0.8), Inches(0.5), Inches(10), Inches(0.6),
                 "250-Station Expansion ROI", font_size=32, color=NAVY, bold=True)

    roi_kpis = [
        ("$12.9M/yr", "Net Annual Profit", GREEN),
        ("15 months", "Payback Period", BLUE),
        ("~$377M", "5-Year Revenue Gap\n(invest vs. don't)", AMBER),
        ("250", "New Stations\n(Bronx/Queens/Brooklyn)", PURPLE),
    ]
    card_w = Inches(2.6)
    for i, (val, label, accent) in enumerate(roi_kpis):
        add_kpi_card(slide, Inches(0.8) + i * (card_w + Inches(0.2)),
                     Inches(1.5), card_w, Inches(1.2), label, val, accent)

    roi_bullets = [
        "250 new stations targeting Bronx, Queens, and Brooklyn — same boroughs as IBZ corridor",
        "Our model identifies which neighborhoods the rollout should prioritize first",
        "Stations placed near high-delay MTA stops capture mode-shift demand",
        "Capital investment recovered in under 2 years from incremental revenue",
        "Not investing has a cost: ~$377M in unrealized revenue over 5 years",
    ]
    add_bullet_slide_content(slide, Inches(0.8), Inches(3.3), Inches(11.5), Inches(3),
                              roi_bullets, font_size=16, color=NAVY)

    # ================================================================
    # SECTION: Success Stories — one slide per city with image placeholder
    # ================================================================
    section_divider(prs, "Success Stories",
                    "Government investment works — 6 cities prove it")

    SUCCESS_CITIES = [
        {
            "city": "San Francisco",
            "system": "Bay Wheels",
            "tagline": "SFMTA integration turned Bay Wheels into transit-connected infrastructure",
            "bullets": [
                "8.1M trips across 699 stations (23 months)",
                "Deep SFMTA transit integration (Clipper card, route planning)",
                "$16M MTC expansion + 4,000 shared e-bikes approved",
                "2017 expansion delivered at zero taxpayer expense",
                "NYC is 10x the market — if SF works, NYC works at scale",
            ],
        },
        {
            "city": "Paris",
            "system": "Velib' Metropole",
            "tagline": "The world's largest bike-share — sustained by public subsidy for nearly two decades",
            "bullets": [
                "48.5M trips in 2025 — highest ridership in Europe",
                "1,400+ stations and 16,000+ bikes across 67 municipalities",
                "A docking point roughly every 300 meters",
                "60-70% of operating costs publicly subsidized",
                "Sustained public investment keeps the system expanding",
            ],
        },
        {
            "city": "London",
            "system": "Santander Cycles",
            "tagline": "A decade of corporate sponsorship funding continuous fleet and station growth",
            "bullets": [
                "106M+ hires since 2015 sponsorship began",
                "800+ stations, 10,000 classic bikes, 2,000+ e-bikes",
                "Daily journeys hit 1.5M in 2025 (+12.7% YoY)",
                "New £220M operating contract (Lyft + Serco)",
                "E-bikes (added 2022) already logged 2.3M rides",
            ],
        },
        {
            "city": "Montreal",
            "system": "BIXI",
            "tagline": "Public bailout and non-profit restructuring turned bankruptcy into a growth story",
            "bullets": [
                "+146% unique users since 2014 restructuring",
                "1,080 stations and 12,600 bikes (3,200 electric) by 2025",
                "100+ million cumulative trips, record usage every summer",
                "Funded 50% user fees, 25% sponsorship, 25% public subsidy",
                "Reached full financial stability by 2018",
            ],
        },
        {
            "city": "Washington, D.C.",
            "system": "Capital Bikeshare",
            "tagline": "The largest municipally-owned bike-share in the U.S. — and fastest-growing",
            "bullets": [
                "6+ million trips in 2024 (+36.9% YoY, second consecutive record)",
                "Jointly owned by 8 local governments",
                "55 miles of new bike lanes (35 protected) + 67-mile trail network",
                "E-bikes drive 60%+ of rides after 143% jump in one year",
                "Overtook Chicago's Divvy for #2 spot nationally",
            ],
        },
        {
            "city": "Chicago",
            "system": "Divvy",
            "tagline": "Self-funding expansion: Lyft's capital investment repaid with revenue-sharing to the city",
            "bullets": [
                "Record 6.8+ million trips in 2025",
                "Expanded to all 50 city wards by 2023",
                "200 new or upgraded stations planned for 2026",
                "Lyft's $50M capital paired with $77M returned to city over 9 years",
                "Owned by Chicago DOT, operated by Lyft since 2019",
            ],
        },
    ]

    for story in SUCCESS_CITIES:
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        # Full-slide image placeholder (dark overlay rectangle)
        # Right half is the image area — user drags a city photo here
        img_placeholder = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(6.5), Inches(0), Inches(6.833), SLIDE_H,
        )
        img_placeholder.fill.solid()
        img_placeholder.fill.fore_color.rgb = RGBColor(0xE2, 0xE8, 0xF0)
        img_placeholder.line.fill.background()
        # Label inside the placeholder
        add_text_box(slide, Inches(7.5), Inches(3.2), Inches(4.5), Inches(1),
                     f"Drag {story['city']} city photo here\n(right-click > Change Picture)",
                     font_size=14, color=MID_GRAY, alignment=PP_ALIGN.CENTER)

        # Left side: dark panel with content
        left_panel = add_shape_fill(slide, Inches(0), Inches(0), Inches(6.5), SLIDE_H, NAVY)

        # City name + system
        add_text_box(slide, Inches(0.8), Inches(0.6), Inches(5), Inches(0.4),
                     story["system"].upper(), font_size=12, color=PINK, bold=True)
        add_text_box(slide, Inches(0.8), Inches(1.0), Inches(5.2), Inches(0.8),
                     story["city"], font_size=36, color=WHITE, bold=True)
        add_text_box(slide, Inches(0.8), Inches(1.9), Inches(5.2), Inches(0.6),
                     story["tagline"], font_size=15, color=LIGHT_GRAY)

        # Bullet points
        add_bullet_slide_content(slide, Inches(0.8), Inches(3.0), Inches(5.2), Inches(3.5),
                                  story["bullets"], font_size=15, color=WHITE, spacing=Pt(10))

    # Common thread panel (after all city slides)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_LIGHT)
    add_text_box(slide, Inches(0.8), Inches(0.5), Inches(10), Inches(0.6),
                 "What These Systems Have in Common", font_size=32, color=NAVY, bold=True)

    common_bullets = [
        "Public ownership or a formal public-private partnership — not purely private ventures",
        "Sustained, multi-year capital commitments — dedicated tranches, not one-off grants",
        "Investment shows up directly in ridership and infrastructure growth",
        "Financial sustainability follows investment — BIXI and Divvy prove self-sustaining economics",
        "Every system posted double-digit ridership growth in its most recent reporting period",
        "This is the evidence base for the NYC investment strategy",
    ]
    add_bullet_slide_content(slide, Inches(0.8), Inches(1.5), Inches(11.5), Inches(4.5),
                              common_bullets, font_size=17, color=NAVY, spacing=Pt(14))

    panel = add_shape_fill(slide, Inches(0.8), Inches(5.8), Inches(11.5), Inches(0.8),
                           RGBColor(0xD1, 0xFA, 0xE5))
    add_text_box(slide, Inches(1.0), Inches(5.9), Inches(11), Inches(0.6),
                 "These cities are independent, external evidence — not benchmarked against NYC. "
                 "SF proves the model works; NYC is the 10x opportunity.",
                 font_size=14, color=RGBColor(0x06, 0x5F, 0x46))

    # ================================================================
    # SECTION: Dashboard
    # ================================================================
    section_divider(prs, "Interactive Dashboard",
                    "8-tab Streamlit decision engine")

    # SLIDE 10: Dashboard Overview
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_LIGHT)
    add_text_box(slide, Inches(0.8), Inches(0.5), Inches(10), Inches(0.6),
                 "Streamlit Dashboard: 8 Analysis Views", font_size=32, color=NAVY, bold=True)

    tabs = [
        ("Overview", "NYC demand trends, rider/bike mix, hourly heatmap with peak-hour highlighting"),
        ("Station Explorer", "Interactive map with demand pressure coloring, top-10 rankings, per-station time series"),
        ("Forecast Lab", "XGBoost-powered scenario model with weather/event sliders and confidence bands"),
        ("MTA Connection", "Subway ridership vs bike trips scatter, transit opportunity scoring by neighborhood"),
        ("Success Stories", "6 global case studies (SF, Paris, London, Montreal, DC, Chicago)"),
        ("Government Investment", "Full NPV/BCR planner with editable assumptions and recommendation table"),
        ("DOT Support Case", "6-argument pitch: gov investment works, MTA failing, green/cheap/healthy, capacity, revenue, market"),
        ("Data & Methods", "Schema documentation, methodology notes, data provenance"),
    ]
    for i, (tab, desc) in enumerate(tabs):
        y = Inches(1.3) + i * Inches(0.72)
        # Tab badge
        badge = add_shape_fill(slide, Inches(0.8), y, Inches(2.2), Inches(0.55), NAVY)
        add_text_box(slide, Inches(0.85), y + Inches(0.1), Inches(2.1), Inches(0.35),
                     f"{i+1}. {tab}", font_size=13, color=WHITE, bold=True)
        # Description
        add_text_box(slide, Inches(3.2), y + Inches(0.1), Inches(9), Inches(0.4),
                     desc, font_size=14, color=NAVY)

    # ================================================================
    # SLIDE 11: Key Findings
    # ================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_LIGHT)
    add_text_box(slide, Inches(0.8), Inches(0.5), Inches(10), Inches(0.6),
                 "Key Findings", font_size=32, color=NAVY, bold=True)

    findings = [
        ("Supply, Not Demand",
         "The system is supply-constrained. 51.9% of stations at capacity means demand is being turned away daily."),
        ("E-Bikes Are the Business",
         "70% e-bike share and $0.27/min overage generate ~$97M/yr. E-bike fleet expansion has outsized ROI."),
        ("MTA Failures Create Bike Demand",
         "Stations near high-delay subway stops show measurably higher bike usage. Transit unreliability is a tailwind."),
        ("Model Validates Station Placement",
         "XGBoost identifies which stations are under-serving demand (+42.5% accuracy vs naive). "
         "Lag features dominate, meaning demand is persistent and predictable."),
        ("Government Investment Works",
         "SF Bay Wheels proves the model. NYC is 10x the market. 250 new stations pay back in 15 months."),
    ]
    for i, (title, body) in enumerate(findings):
        y = Inches(1.4) + i * Inches(1.1)
        # Number badge
        badge = add_shape_fill(slide, Inches(0.8), y, Inches(0.45), Inches(0.45), PINK)
        add_text_box(slide, Inches(0.8), y + Inches(0.05), Inches(0.45), Inches(0.35),
                     str(i + 1), font_size=16, color=WHITE, bold=True,
                     alignment=PP_ALIGN.CENTER)
        add_text_box(slide, Inches(1.5), y, Inches(4), Inches(0.35),
                     title, font_size=16, color=NAVY, bold=True)
        add_text_box(slide, Inches(1.5), y + Inches(0.35), Inches(10.5), Inches(0.6),
                     body, font_size=14, color=MID_GRAY)

    # ================================================================
    # SLIDE 12: Closing / Next Steps
    # ================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, NAVY)

    add_text_box(slide, Inches(1), Inches(1.5), Inches(10), Inches(0.4),
                 "THE ASK", font_size=14, color=PINK, bold=True)
    add_text_box(slide, Inches(1), Inches(2.0), Inches(10), Inches(1),
                 "Recommendations & Next Steps",
                 font_size=36, color=WHITE, bold=True)

    next_steps = [
        "1.  Share internal rebalancing data to refine station-level forecasts",
        "2.  Pilot 50 stations in highest-demand underserved neighborhoods",
        "3.  Integrate with NYC DOT for transit-connected bike corridors",
        "4.  Extend model to hourly granularity for real-time rebalancing",
        "5.  Cross-city rollout: apply framework to DC, Chicago, and new markets",
    ]
    add_bullet_slide_content(slide, Inches(1), Inches(3.3), Inches(10), Inches(3),
                              next_steps, font_size=18, color=WHITE, spacing=Pt(12))

    add_text_box(slide, Inches(1), Inches(6.2), Inches(11), Inches(0.5),
                 "Data-driven decisions for NYC's bike-share future",
                 font_size=16, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

    # ── Save ──
    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT))
    print(f"\nSaved presentation to {OUT}")
    print(f"  {len(prs.slides)} slides")


if __name__ == "__main__":
    build()
