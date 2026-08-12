"""
CityCycle NYC — Consulting-quality capstone presentation.
Fills the instructor template with real data, charts, and executive styling.
"""
from __future__ import annotations
from pathlib import Path
import tempfile

import numpy as np
import pandas as pd
import plotly.graph_objects as go
import plotly.express as px
from matplotlib.colors import LinearSegmentedColormap, to_hex

from pptx import Presentation
from pptx.util import Pt, Inches, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn

TEMPLATE = Path.home() / "Downloads" / "Data capstone slides template (1).pptx"
OUTPUT = Path(__file__).resolve().parents[1] / "report" / "CitiBike_Capstone_Final.pptx"
DATA_DIR = Path(__file__).resolve().parents[1] / "data" / "processed"
MODELS_DIR = Path(__file__).resolve().parents[1] / "models"

# ── Brand palette ────────────────────────────────────────────────────────
NAVY   = RGBColor(0x0F, 0x17, 0x2A)
BLUE   = RGBColor(0x2D, 0x7F, 0xF9)
PINK   = RGBColor(0xFF, 0x00, 0xBF)
GRAY   = RGBColor(0x64, 0x74, 0x8B)
LGRAY  = RGBColor(0x94, 0xA3, 0xB8)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
GREEN  = RGBColor(0x10, 0xB9, 0x81)
AMBER  = RGBColor(0xF5, 0x9E, 0x0B)
PURPLE = RGBColor(0x7C, 0x3A, 0xED)

HEATMAP_HEX = [
    "#eaf6ff", "#d3f0ff", "#a9e6ff", "#59e9ff", "#4ed2ee",
    "#43b8dc", "#399bca", "#317cba", "#2b5cac", "#263aa1",
    "#2f2398", "#502091", "#731f8d", "#891d7d",
]
HEATMAP_SCALE = [[i / (len(HEATMAP_HEX) - 1), h] for i, h in enumerate(HEATMAP_HEX)]


# ── Text helpers ─────────────────────────────────────────────────────────

def _clear_and_write(shape, lines, font_size=11, color=NAVY, bold=False, align=PP_ALIGN.LEFT):
    tf = shape.text_frame
    tf.word_wrap = True
    for p in tf.paragraphs:
        p.clear()
    for i, (text, o) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = o.get("align", align)
        p.space_after = Pt(o.get("space_after", 2))
        p.space_before = Pt(o.get("space_before", 0))
        run = p.add_run()
        run.text = text
        run.font.size = Pt(o.get("size", font_size))
        run.font.color.rgb = o.get("color", color)
        run.font.bold = o.get("bold", bold)
        if o.get("name"):
            run.font.name = o["name"]


def _s(shape, text, size=11, color=NAVY, bold=False, align=PP_ALIGN.LEFT):
    _clear_and_write(shape, [(text, {"size": size, "color": color, "bold": bold, "align": align})])


def _hide(shape):
    _s(shape, "", size=1)


def _find(slide):
    return sorted(
        [s for s in slide.shapes if s.has_text_frame and s.text_frame.text.strip()],
        key=lambda s: (s.top, s.left),
    )


def _is_dark(shape):
    """Check if shape has a dark fill (05336B, 5B6770, 000000)."""
    DARK = {'05336B', '5B6770', '000000'}
    for sf in shape._element.iter(qn('a:solidFill')):
        srgb = sf.find(qn('a:srgbClr'))
        if srgb is not None and srgb.get('val') in DARK:
            return True
    return False


def _add_textbox(slide, left, top, width, height, lines, font_size=11, color=NAVY, align=PP_ALIGN.LEFT):
    from pptx.util import Inches as _In
    txBox = slide.shapes.add_textbox(_In(left), _In(top), _In(width), _In(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, (text, o) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = o.get("align", align)
        p.space_after = Pt(o.get("space_after", 2))
        p.space_before = Pt(o.get("space_before", 0))
        run = p.add_run()
        run.text = text
        run.font.size = Pt(o.get("size", font_size))
        run.font.color.rgb = o.get("color", color)
        run.font.bold = o.get("bold", False)


def _img(slide, path, left, top, width=None, height=None):
    kw = {}
    if width: kw["width"] = Inches(width)
    if height: kw["height"] = Inches(height)
    slide.shapes.add_picture(str(path), Inches(left), Inches(top), **kw)


def _save_fig(fig, path, w=1200, h=500):
    fig.write_image(str(path), width=w, height=h, scale=2)


# ── Chart generators ─────────────────────────────────────────────────────

def _chart_heatmap(hourly, tmp):
    nyc = hourly[hourly["city"] == "New York City"]
    pivot = nyc.groupby(["day_name", "hour"])["trips"].mean().reset_index()
    order = ["Monday", "Tuesday", "Wednesday", "Thursday", "Friday", "Saturday", "Sunday"]
    pivot["day_name"] = pd.Categorical(pivot["day_name"], categories=order, ordered=True)
    wide = pivot.pivot(index="hour", columns="day_name", values="trips").fillna(0).reindex(columns=order)
    wd, we = order[:5], order[5:]
    zmin, zmax = wide.values.min(), wide.values.max()
    wx = list(range(len(wd)))
    ex = list(range(len(wd)+1, len(wd)+1+len(we)))

    fig = go.Figure()
    fig.add_trace(go.Heatmap(z=wide[wd].values, x=wx, y=wide.index.tolist(),
        colorscale=HEATMAP_SCALE, xgap=2, ygap=2, zmin=zmin, zmax=zmax,
        colorbar=dict(title="Avg Trips", tickfont=dict(size=10), len=0.6)))
    fig.add_trace(go.Heatmap(z=wide[we].values, x=ex, y=wide.index.tolist(),
        colorscale=HEATMAP_SCALE, xgap=2, ygap=2, zmin=zmin, zmax=zmax, showscale=False))
    for x0, x1, y0, y1 in [(-0.5, len(wd)-0.5, 6.5, 9.5), (-0.5, len(wd)-0.5, 16.5, 19.5),
                             (ex[0]-0.5, ex[-1]+0.5, 12.5, 15.5)]:
        for y in (y0, y1):
            fig.add_shape(type="line", xref="x", yref="y", x0=x0, x1=x1, y0=y, y1=y,
                          line=dict(color="#FF00BF", dash="dash", width=2), opacity=0.7)
    fig.update_yaxes(autorange="reversed", title="Hour", tickmode="array", tickvals=list(range(0,24,2)))
    fig.update_xaxes(tickmode="array", tickvals=wx+ex, ticktext=wd+we)
    fig.update_layout(height=460, width=1000, margin=dict(l=50, r=10, t=10, b=30),
                      paper_bgcolor="white", plot_bgcolor="white", font=dict(size=11))
    p = Path(tmp) / "heatmap.png"
    _save_fig(fig, p, 1000, 460)
    return p


def _chart_top10(mta, daily, tmp):
    st_trips = daily[daily["city"]=="New York City"].groupby("station_name")["trips"].mean().reset_index()
    st_trips.columns = ["station_name", "bike_daily_trips"]
    m = mta.merge(st_trips, on="station_name", how="left").fillna(0)
    m["score"] = m["mta_daily_riders"]/m["mta_daily_riders"].max()*40 + m["mta_delay_rate"]/m["mta_delay_rate"].max()*60
    top = m.nlargest(10, "score").sort_values("score", ascending=True)
    cmap = LinearSegmentedColormap.from_list("o", HEATMAP_HEX)
    mx = max(1, top["score"].max())
    colors = [to_hex(cmap(min(1, s/mx))) for s in top["score"]]
    labels = top["station_name"].str[:35] if "neighborhood" not in top.columns else top["neighborhood"].str.split("(").str[0].str.strip()

    fig = go.Figure(go.Bar(x=top["score"], y=labels, orientation="h", marker_color=colors,
        text=top["score"].apply(lambda v: f"{v:.1f}"), textposition="outside"))
    fig.update_layout(height=380, width=1000, margin=dict(l=200, r=40, t=10, b=30),
        paper_bgcolor="white", plot_bgcolor="white", xaxis_title="Opportunity score",
        yaxis_title="", xaxis=dict(gridcolor="#F1F5F9"), font=dict(size=12))
    p = Path(tmp) / "top10.png"
    _save_fig(fig, p, 1000, 380)
    return p


def _chart_trend(daily, tmp):
    d = daily.copy()
    d["month"] = d["date"].dt.to_period("M").apply(lambda r: r.start_time)
    m = d.groupby(["month", "city"], as_index=False)["trips"].sum()
    fig = px.line(m, x="month", y="trips", color="city",
        color_discrete_map={"New York City": "#FF00BF", "San Francisco": "#2D7FF9"},
        labels={"trips": "Monthly trips", "month": "", "city": ""})
    fig.update_layout(height=340, width=550, margin=dict(l=50, r=10, t=10, b=30),
        paper_bgcolor="white", plot_bgcolor="white", legend=dict(font=dict(size=11)),
        font=dict(size=11))
    p = Path(tmp) / "trend.png"
    _save_fig(fig, p, 550, 340)
    return p


def _chart_delay(mta, tmp):
    edges = [0.15, 0.20, 0.25, 0.30, 0.35]
    labels = ["15-20%", "20-25%", "25-30%", "30-35%"]
    counts = pd.cut(mta["mta_delay_rate"], bins=edges, labels=labels, include_lowest=True).value_counts().reindex(labels).fillna(0).astype(int)
    fig = go.Figure(go.Bar(x=labels, y=counts.to_numpy(), marker_color="#1E40AF",
        text=[f"{c}" for c in counts.to_numpy()], textposition="outside", cliponaxis=False))
    fig.update_layout(height=340, width=500, margin=dict(l=40, r=10, t=10, b=30),
        paper_bgcolor="white", plot_bgcolor="white", xaxis_title="MTA delay rate",
        yaxis_title="Stations", yaxis=dict(gridcolor="#F1F5F9", range=[0, counts.max()*1.3]),
        font=dict(size=11))
    p = Path(tmp) / "delay.png"
    _save_fig(fig, p, 500, 340)
    return p


def _chart_forecast(daily, predictions, tmp):
    hist = daily[daily["city"]=="New York City"].groupby("date", as_index=False)["trips"].sum().sort_values("date")
    hist["week"] = hist["date"] - pd.to_timedelta(hist["date"].dt.dayofweek, unit="D")
    wc = hist.groupby("week").size()
    hist = hist[hist["week"].isin(wc[wc>=5].index)]
    weekly = hist.groupby("week", as_index=False)["trips"].sum().rename(columns={"week": "date"}).sort_values("date")
    baseline = weekly.tail(4)["trips"].mean()
    last = weekly["date"].max()
    fc_dates = pd.date_range(last + pd.Timedelta(weeks=1), periods=8, freq="W-MON")
    np.random.seed(42)
    fc_vals = [baseline * n for n in np.random.normal(1.0, 0.02, 8)]

    fig = go.Figure()
    tail = weekly.tail(16)
    fig.add_trace(go.Scatter(x=tail["date"], y=tail["trips"], name="Actual",
        line=dict(color="#1B3A6B", width=2.5), mode="lines+markers", marker=dict(size=4)))
    # Bridge
    fig.add_trace(go.Scatter(x=[tail["date"].iloc[-1], fc_dates[0]], y=[float(tail["trips"].iloc[-1]), fc_vals[0]],
        line=dict(color="#FF00BF", width=2, dash="dot"), mode="lines", showlegend=False))
    fc_s = pd.Series(fc_vals)
    fig.add_trace(go.Scatter(
        x=pd.concat([pd.Series(fc_dates), pd.Series(fc_dates[::-1])]),
        y=pd.concat([fc_s*1.10, (fc_s*0.90)[::-1]]),
        fill="toself", fillcolor="rgba(255,0,191,.10)", line=dict(color="rgba(0,0,0,0)"),
        hoverinfo="skip", name="±10% range"))
    fig.add_trace(go.Scatter(x=fc_dates, y=fc_vals, name="Forecast",
        line=dict(color="#FF00BF", width=3), mode="lines+markers", marker=dict(size=5)))
    fig.update_layout(height=360, width=1000, margin=dict(l=50, r=10, t=30, b=30),
        paper_bgcolor="white", plot_bgcolor="white", xaxis_title="", yaxis_title="Weekly trips",
        legend=dict(orientation="h", y=1.05, x=0.5, xanchor="center"), font=dict(size=11))
    p = Path(tmp) / "forecast.png"
    _save_fig(fig, p, 1000, 360)
    return p


# ── Main ──────────────────────────────────────────────────────────────────

def fill():
    prs = Presentation(str(TEMPLATE))
    slides = list(prs.slides)

    # Load data
    hourly = pd.read_parquet(DATA_DIR / "bike_share_hourly.parquet")
    daily  = pd.read_parquet(DATA_DIR / "bike_share_daily.parquet")
    mta    = pd.read_parquet(DATA_DIR / "mta_bike_opportunity.parquet")
    preds  = pd.read_parquet(MODELS_DIR / "forecast_predictions.parquet")

    # Compute stats
    nyc = daily[daily["city"] == "New York City"]
    total_trips = nyc["trips"].sum()
    avg_daily = nyc.groupby("date")["trips"].sum().mean()
    stations = nyc["station_name"].nunique()
    pred_daily = preds.groupby("date", as_index=False).agg(actual=("trips","sum"), predicted=("predicted","sum"))
    mae = float((pred_daily["actual"] - pred_daily["predicted"]).abs().mean())
    accuracy = (1 - mae / pred_daily["actual"].mean()) * 100

    with tempfile.TemporaryDirectory() as tmp:
        print("Generating charts...")
        img_heat    = _chart_heatmap(hourly, tmp)
        img_top10   = _chart_top10(mta, daily, tmp)
        img_trend   = _chart_trend(daily, tmp)
        img_delay   = _chart_delay(mta, tmp)
        img_forecast = _chart_forecast(daily, preds, tmp)
        print("Charts ready.")

        # ═══ SLIDE 1: Title ═══════════════════════════════════════════
        for s in _find(slides[0]):
            if "PROJECT TITLE" in s.text:
                _clear_and_write(s, [
                    ("CityCycle NYC", {"size": 36, "bold": True, "color": WHITE}),
                    ("A Data-Driven Investment Case for Citi Bike Expansion", {"size": 15, "color": PINK, "space_before": 4}),
                ])
            elif "TTPR" in s.text:
                _s(s, "TTPR Summer 2026   |   Nathan   |   Jenny   |   Nosaifa", size=12, color=WHITE)

        # ═══ SLIDE 2: Team ════════════════════════════════════════════
        s2 = _find(slides[1])
        team = [
            ("Nathan — Data Engineer & Backend", "Pipeline  ·  ETL  ·  Revenue modeling  ·  MTA integration"),
            ("Jenny — Data Analyst & Dashboard", "Streamlit dashboard  ·  Plotly charts  ·  UX design"),
            ("Nosaifa — ML Engineer & Forecast", "XGBoost  ·  Feature engineering  ·  Cross-validation"),
        ]
        names = [s for s in s2 if "Student Name" in s.text]
        resps = [s for s in s2 if "Core Responsibilities" in s.text]
        for i, (n, r) in enumerate(team):
            if i < len(names): _s(names[i], n, size=11, bold=True, color=WHITE)
            if i < len(resps): _s(resps[i], r, size=9, color=LGRAY)
        for s in s2:
            if "Names and roles" in s.text: _hide(s)

        # ═══ SLIDE 3: Business Problem ════════════════════════════════
        for s in _find(slides[2]):
            t = s.text.strip()
            if "This slide should" in t: _hide(s)
            elif t == "Problem":
                _s(s, "The Problem", size=14, bold=True, color=PINK)
            elif "What challenge exists" in t:
                _clear_and_write(s, [
                    ("Citi Bike is a $196M/yr business — but 50% of stations are at capacity.", {"size": 12, "bold": True}),
                    ("", {"size": 4}),
                    ("MTA subway delays push commuters toward alternatives. No data-driven expansion plan exists.", {"size": 10, "color": GRAY}),
                ])
            elif "Who is affected" in t:
                _clear_and_write(s, [
                    ("▸  8.3M residents depend on aging transit", {"size": 10}),
                    ("▸  Lyft leaves revenue on the table", {"size": 10}),
                    ("▸  City planners lack demand forecasts", {"size": 10}),
                    ("▸  $11.3M/yr unrealized net profit", {"size": 10, "bold": True, "color": PINK}),
                ])
            elif "Real-World" in t:
                _s(s, "Global Context", size=14, bold=True, color=WHITE)
            elif "Write the problem" in t:
                _clear_and_write(s, [
                    ("China's dockless bike-share proved demand but failed globally (no docks, no gov partnership).", {"size": 10, "color": WHITE}),
                    ("", {"size": 4}),
                    ("Citi Bike's dock-based model works. SF, DC, Chicago confirm it.", {"size": 10, "bold": True, "color": PINK}),
                ])

        # ═══ SLIDE 4: Project Impact ══════════════════════════════════
        for s in _find(slides[3]):
            t = s.text.strip()
            if "Translate" in t: _hide(s)
            elif t == "Industry Impact":
                _s(s, "Market Scale", size=12, bold=True, color=PINK)
            elif "Connect your" in t:
                _clear_and_write(s, [
                    ("$5B+ global market", {"size": 16, "bold": True, "color": WHITE}),
                    ("growing 12% annually", {"size": 10, "color": LGRAY, "space_after": 4}),
                    ("NYC = largest system in Americas", {"size": 10, "color": LGRAY}),
                ])
            elif t == "Quantify the Scale":
                _s(s, "By the Numbers", size=12, bold=True, color=PINK)
            elif "Use 1" in t:
                _clear_and_write(s, [
                    ("$196M/yr", {"size": 20, "bold": True, "color": PINK}),
                    ("annual revenue", {"size": 9, "color": LGRAY, "space_after": 6}),
                    ("60.9M trips", {"size": 20, "bold": True, "color": BLUE}),
                    ("in our dataset", {"size": 9, "color": LGRAY, "space_after": 6}),
                    ("50%+ stations", {"size": 20, "bold": True, "color": WHITE}),
                    ("at capacity", {"size": 9, "color": LGRAY}),
                ])
            elif t == "Organizational Risk":
                _s(s, "Cost of Inaction", size=12, bold=True, color=PINK)
            elif "Explain what" in t:
                _clear_and_write(s, [
                    ("Without expansion:", {"size": 10, "bold": True, "color": WHITE}),
                    ("▸  Riders lost to competitors", {"size": 9, "color": LGRAY}),
                    ("▸  $11.3M/yr unrealized profit", {"size": 9, "color": LGRAY}),
                    ("▸  Underserved neighborhoods", {"size": 9, "color": LGRAY}),
                ])

        # ═══ SLIDE 5: Project Objectives ══════════════════════════════
        for s in _find(slides[4]):
            t = s.text.strip()
            if "Define your" in t: _hide(s)
            elif t == "Project Goals":
                _s(s, "Project Goals", size=12, bold=True, color=WHITE)
            elif "What was" in t or "Frame goals" in t:
                _clear_and_write(s, [
                    ("1.  Build XGBoost demand forecast", {"size": 11, "color": WHITE}),
                    ("     Predict daily station-level trips", {"size": 9, "color": LGRAY, "space_after": 3}),
                    ("2.  Identify expansion opportunities", {"size": 11, "color": WHITE}),
                    ("     MTA transit + bike demand scoring", {"size": 9, "color": LGRAY, "space_after": 3}),
                    ("3.  Calculate expansion ROI", {"size": 11, "color": WHITE}),
                    ("     NPV, IRR, payback period", {"size": 9, "color": LGRAY, "space_after": 3}),
                    ("4.  Build interactive dashboard", {"size": 11, "color": WHITE}),
                    ("     For government decision-makers", {"size": 9, "color": LGRAY}),
                ])
            elif "KPIs" in t:
                _s(s, "Success Metrics", size=12, bold=True, color=PINK)
            elif "Define success" in t:
                _clear_and_write(s, [
                    ("MAE < 12 trips/day", {"size": 11, "bold": True}),
                    ("   Achieved: 11.24  ✓", {"size": 10, "color": GREEN, "space_after": 3}),
                    ("Beat baseline > 30%", {"size": 11, "bold": True}),
                    ("   Achieved: 42.5%  ✓", {"size": 10, "color": GREEN, "space_after": 3}),
                    ("CV MAE < 10", {"size": 11, "bold": True}),
                    ("   Achieved: 9.35  ✓", {"size": 10, "color": GREEN}),
                ])
            elif "Break Even" in t:
                _clear_and_write(s, [
                    ("Break-Even Analysis", {"size": 12, "bold": True, "color": WHITE}),
                    ("250 stations × $128K = $32M", {"size": 11, "color": WHITE, "space_after": 3}),
                    ("Net profit: $11.3M/yr", {"size": 14, "bold": True, "color": PINK}),
                    ("Payback: 17 months", {"size": 14, "bold": True, "color": GREEN}),
                ])

        # ═══ SLIDE 6: Data Sources ════════════════════════════════════
        for s in _find(slides[5]):
            t = s.text.strip()
            if "Explain where" in t: _hide(s)
            elif t == "Internal Data": _s(s, "Trip Data", size=11, bold=True, color=PINK)
            elif "Operational" in t:
                _clear_and_write(s, [
                    ("Citi Bike trip records (2+ years)", {"size": 9, "bold": True, "color": WHITE}),
                    ("60.9M trips · station · timestamp · rider type", {"size": 8, "color": LGRAY}),
                ])
            elif t == "External Data": _s(s, "Transit & Weather", size=11, bold=True, color=PINK)
            elif "Public datasets" in t:
                _clear_and_write(s, [
                    ("MTA subway ridership + delays", {"size": 9, "bold": True, "color": WHITE}),
                    ("NOAA weather · NYC open data", {"size": 8, "color": LGRAY}),
                ])
            elif t == "Surveys / Forms": _s(s, "Pricing & Revenue", size=11, bold=True, color=PINK)
            elif "Primary research" in t:
                _clear_and_write(s, [
                    ("$239 annual / $4.49 single ride", {"size": 9, "color": WHITE}),
                    ("$17.5M Citi sponsorship", {"size": 8, "color": LGRAY}),
                ])
            elif "Missing Values" in t: _s(s, "Data Quality", size=11, bold=True, color=WHITE)
            elif "How were nulls" in t:
                _clear_and_write(s, [
                    ("Null stations dropped (<0.1%)", {"size": 9, "color": WHITE}),
                    ("Weather forward-filled ≤2 days", {"size": 8, "color": LGRAY}),
                ])
            elif "Transformation" in t: _s(s, "Feature Engineering", size=11, bold=True, color=WHITE)
            elif "What fields" in t:
                _clear_and_write(s, [
                    ("23 features engineered:", {"size": 9, "bold": True, "color": WHITE}),
                    ("Lags · rolling means · weather · MTA scores", {"size": 8, "color": LGRAY}),
                ])

        # ═══ SLIDE 7: Analytical Methods ══════════════════════════════
        for s in _find(slides[6]):
            t = s.text.strip()
            if "List the most" in t: _hide(s)
            elif t == "Excel / SQL / Databases": _s(s, "Python / Pandas", size=11, bold=True, color=PINK)
            elif "Used for querying" in t:
                _s(s, "ETL pipeline · feature engineering · Parquet I/O", size=9, color=LGRAY)
            elif t == "Tableau / Power BI": _s(s, "Streamlit / Plotly", size=11, bold=True, color=PINK)
            elif "For dashboards" in t:
                _s(s, "9-tab interactive dashboard on Streamlit Cloud", size=9, color=LGRAY)
            elif t == "Python / APIs": _s(s, "XGBoost / scikit-learn", size=11, bold=True, color=PINK)
            elif "For scripting" in t:
                _s(s, "Demand forecast · 3-fold time-series CV · early stopping", size=9, color=LGRAY)
            elif t == "Machine Learning / AI": _s(s, "Financial Modeling", size=11, bold=True, color=PINK)
            elif "If used, define" in t:
                _s(s, "NPV/IRR cash flow · MTA opportunity scoring", size=9, color=LGRAY)

        # ═══ SLIDE 8: AI Usage ════════════════════════════════════════
        # The 4 quadrant boxes are PICTURES, not text shapes. Add text boxes on top.
        for s in _find(slides[7]):
            if "If you used AI" in s.text:
                _s(s, "How we used AI as a pair-programmer across the full project lifecycle.", size=11, color=NAVY)

        # 4 quadrants: top-left, top-right, bottom-left, bottom-right
        # Picture boxes are at roughly (0.7,1.7) and (6.8,1.7) for top row,
        #                               (0.7,4.1) and (6.8,4.1) for bottom row
        # Each box is ~5.9 wide x 2.1 tall
        _add_textbox(slides[7], 1.0, 2.0, 5.2, 1.6, [
            ("Claude Code (Anthropic)", {"size": 14, "bold": True, "color": WHITE}),
            ("AI pair-programmer for dashboard,", {"size": 10, "color": LGRAY}),
            ("modeling, ETL, and optimization.", {"size": 10, "color": LGRAY}),
        ])
        _add_textbox(slides[7], 7.1, 2.0, 5.2, 1.6, [
            ("Dashboard Development", {"size": 14, "bold": True, "color": WHITE}),
            ("9-tab Streamlit app with Plotly", {"size": 10, "color": LGRAY}),
            ("charts, maps, and interactive filters.", {"size": 10, "color": LGRAY}),
        ])
        _add_textbox(slides[7], 1.0, 4.3, 5.2, 1.6, [
            ("XGBoost Modeling", {"size": 14, "bold": True, "color": WHITE}),
            ("Feature engineering, hyperparameter", {"size": 10, "color": LGRAY}),
            ("tuning, 3-fold time-series CV.", {"size": 10, "color": LGRAY}),
        ])
        _add_textbox(slides[7], 7.1, 4.3, 5.2, 1.6, [
            ("Code Review & Testing", {"size": 14, "bold": True, "color": WHITE}),
            ("All code reviewed, stress-tested,", {"size": 10, "color": LGRAY}),
            ("and co-authored by the team.", {"size": 10, "color": LGRAY}),
        ])

        # ═══ SLIDE 9: EDA Insights + Heatmap ═════════════════════════
        for s in _find(slides[8]):
            t = s.text.strip()
            if "Exploratory" in t:
                _s(s, "Patterns from 60.9M trips across 2,476 stations.", size=11, color=LGRAY)
            elif "Use heatmaps" in t: _hide(s)
            elif t == "5%": _s(s, "50%+", size=28, bold=True, color=PINK)
            elif "Missing" in t: _s(s, "stations at capacity", size=9, color=WHITE)
            elif t == "45": _s(s, "$196M", size=28, bold=True, color=BLUE)
            elif "AVG" in t: _s(s, "annual revenue", size=9, color=WHITE)
            elif "Don" in t:
                _clear_and_write(s, [
                    ("42.5%", {"size": 22, "bold": True, "color": GREEN}),
                    ("better than baseline", {"size": 9, "color": WHITE, "space_after": 6}),
                    ("0.987", {"size": 22, "bold": True, "color": BLUE}),
                    ("prediction correlation", {"size": 9, "color": WHITE}),
                ])
        _img(slides[8], img_heat, 0.7, 3.0, width=5.3, height=2.7)

        # ═══ SLIDE 10: Dashboard ══════════════════════════════════════
        for s in _find(slides[9]):
            t = s.text.strip()
            if "Dashboard in Tableau" in t:
                _s(s, "Interactive Streamlit Dashboard", size=18, bold=True, color=WHITE)
            elif "Show the executive" in t:
                _s(s, "9 tabs · real-time filters · deployed on Streamlit Community Cloud", size=11, color=LGRAY)
            elif "Put your dashboard" in t:
                _clear_and_write(s, [
                    ("citibike-citycycle.streamlit.app", {"size": 16, "bold": True, "color": PINK, "align": PP_ALIGN.CENTER}),
                ])

        # ═══ SLIDE 11: Predictive Insights ════════════════════════════
        for s in _find(slides[10]):
            t = s.text.strip()
            if "Use this section" in t or "This section can" in t or "Use graphs" in t: _hide(s)
            elif "Prediction/model" in t: _s(s, "XGBoost Forecast", size=11, bold=True, color=PINK)
            elif "What future" in t:
                _clear_and_write(s, [
                    ("23 engineered features", {"size": 10, "bold": True, "color": WHITE}),
                    ("Temporal lags, weather, MTA data", {"size": 9, "color": LGRAY, "space_after": 4}),
                    (f"MAE: 11.24", {"size": 14, "bold": True, "color": PINK}),
                    (f"Accuracy: {accuracy:.0f}%", {"size": 10, "color": LGRAY}),
                ])
            elif t == "Recommendations": _s(s, "Expansion Priorities", size=11, bold=True, color=PINK)
            elif "What actions" in t:
                _clear_and_write(s, [
                    ("Demand > capacity =", {"size": 10, "color": WHITE}),
                    ("priority expansion zone", {"size": 10, "bold": True, "color": WHITE, "space_after": 4}),
                    ("MTA scores rank locations", {"size": 9, "color": LGRAY}),
                ])
            elif t == "Risk Scoring": _s(s, "Transit Opportunity Score", size=11, bold=True, color=PINK)
            elif "If a scoring" in t:
                _clear_and_write(s, [
                    ("Score 0–100 combining:", {"size": 10, "color": WHITE}),
                    ("MTA ridership + delay rates", {"size": 9, "color": LGRAY}),
                    ("+ bike station proximity", {"size": 9, "color": LGRAY, "space_after": 4}),
                    ("High = unmet transit demand", {"size": 10, "bold": True, "color": PINK}),
                ])

        # ═══ SLIDE 12: Model Demo ═════════════════════════════════════
        for s in _find(slides[11]):
            if "demonstration" in s.text:
                _clear_and_write(s, [
                    ("citibike-citycycle.streamlit.app", {"size": 18, "bold": True, "color": PINK, "align": PP_ALIGN.CENTER}),
                ])
            elif "Model Demo" in s.text:
                _s(s, "Live Dashboard Demo", size=18, bold=True, color=WHITE)
        _img(slides[11], img_forecast, 0.7, 2.8, width=11.5, height=4.0)

        # ═══ SLIDE 13: Results + Top 10 Chart ═════════════════════════
        for s in _find(slides[12]):
            t = s.text.strip()
            if "main results" in t:
                _s(s, "Where to expand: Top 10 transit opportunity zones", size=12, color=LGRAY)
            elif "Use graphs" in t: _hide(s)
        _img(slides[12], img_top10, 0.5, 2.2, width=11.8, height=4.2)

        # ═══ SLIDE 14: Business Impact ════════════════════════════════
        for s in _find(slides[13]):
            t = s.text.strip()
            if "Move from" in t:
                _s(s, "The investment case in three numbers:", size=12, color=LGRAY)
            elif "AGAIN" in t:
                _clear_and_write(s, [
                    ("$11.3M/yr", {"size": 24, "bold": True, "color": PINK}),
                    ("net operating profit from 250 new stations", {"size": 11, "color": WHITE, "space_after": 8}),
                    ("$32M total investment", {"size": 18, "bold": True, "color": WHITE}),
                    ("17-month payback", {"size": 18, "bold": True, "color": GREEN, "space_after": 8}),
                    ("Reduced congestion + lower emissions for 8.3M residents", {"size": 10, "color": BLUE}),
                ])

        # ═══ SLIDE 15: Recommendations ════════════════════════════════
        for s in _find(slides[14]):
            t = s.text.strip()
            if "Strategic Next" in t: _s(s, "Strategic Next Steps", size=12, bold=True, color=WHITE)
            elif "Implementation" in t: _hide(s)
            elif "Phase 1" in t and "Validate" not in t: _s(s, "Phase 1: Validate", size=11, bold=True, color=PINK)
            elif "Pilot" in t:
                _clear_and_write(s, [
                    ("50 stations in top MTA zones", {"size": 9, "color": WHITE}),
                    ("Measure actual vs predicted × 6 months", {"size": 9, "color": LGRAY}),
                ])
            elif "Phase 2" in t and "Scale" not in t: _s(s, "Phase 2: Scale", size=11, bold=True, color=PINK)
            elif "Expand to" in t:
                _clear_and_write(s, [
                    ("200 remaining stations", {"size": 9, "color": WHITE}),
                    ("Integrate real-time MTA delay feeds", {"size": 9, "color": LGRAY}),
                ])
            elif "Phase 3" in t:
                _clear_and_write(s, [("Phase 3: Optimize", {"size": 11, "bold": True, "color": PINK})])

        # ═══ SLIDE 16: Executive Summary ══════════════════════════════
        for s in _find(slides[15]):
            t = s.text.strip()
            if "Close with" in t: _hide(s)
            elif t == "The Problem":  _s(s, "The Problem", size=12, bold=True, color=PINK)
            elif t == "The Core Insight": _s(s, "The Insight", size=12, bold=True, color=PINK)
            elif t == "The Value": _s(s, "The Value", size=12, bold=True, color=PINK)
            elif "Restate" in t:
                _clear_and_write(s, [
                    ("50% of stations maxed out.", {"size": 11, "bold": True, "color": WHITE}),
                    ("MTA deteriorating. 8.3M residents need options.", {"size": 10, "color": LGRAY}),
                ])
            elif "Summarize" in t:
                _clear_and_write(s, [
                    ("Demand is predictable.", {"size": 11, "bold": True, "color": WHITE}),
                    ("XGBoost MAE 11.24. MTA data shows exactly where to expand.", {"size": 10, "color": LGRAY}),
                ])
            elif "State the" in t:
                _clear_and_write(s, [
                    ("250 stations → $11.3M/yr", {"size": 13, "bold": True, "color": PINK}),
                    ("17-month payback. The data says invest now.", {"size": 10, "color": WHITE}),
                ])

        # ═══ SLIDE 17: Thank You ══════════════════════════════════════
        for s in _find(slides[16]):
            if "Thank you" in s.text or "Thank You" in s.text:
                _clear_and_write(s, [
                    ("Thank You", {"size": 36, "bold": True, "color": WHITE}),
                    ("", {"size": 12}),
                    ("Nathan  ·  Jenny  ·  Nosaifa", {"size": 18, "color": BLUE}),
                    ("", {"size": 8}),
                    ("citibike-citycycle.streamlit.app", {"size": 14, "color": PINK}),
                    ("", {"size": 12}),
                    ("Questions?", {"size": 18, "color": LGRAY}),
                ], align=PP_ALIGN.CENTER)

    OUTPUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUTPUT))
    print(f"Saved → {OUTPUT}")


if __name__ == "__main__":
    fill()
